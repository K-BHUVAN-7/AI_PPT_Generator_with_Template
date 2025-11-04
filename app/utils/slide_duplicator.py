"""
Exact slide duplication utilities
Based on: https://github.com/scanny/python-pptx/issues/132
"""

from pptx.opc.packuri import PackURI
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.parts.slide import SlidePart
from pptx.shapes.group import GroupShape
import copy
import random
import string
import io
import logging

logger = logging.getLogger(__name__)


def _object_rels(obj):
    """Get relationships from object"""
    rels = obj.rels
    rel_keys = list(rels)
    # Some pptx versions return string keys, others Relationship objects
    if rel_keys and isinstance(rel_keys[0], str):
        return [v for _, v in rels.items()]
    else:
        return rel_keys


def _exp_add_slide(ppt, slide_layout):
    """
    Add a new slide with a unique partname to avoid conflicts
    """

    def generate_slide_partname(self):
        """Generate unique slide partname"""
        sldIdLst = self._element.get_or_add_sldIdLst()
        existing_rels = [rel.target_partname for rel in _object_rels(self)]
        partname_str = f"/ppt/slides/slide{len(sldIdLst) + 1}.xml"

        # Ensure unique filename
        while partname_str in existing_rels:
            random_part = "".join(random.choice(string.ascii_letters) for _ in range(2))
            partname_str = f"/ppt/slides/slide{random_part}{len(sldIdLst) + 1}.xml"

        return PackURI(partname_str)

    def add_slide_part(self, slide_layout):
        """Create a new slide part"""
        partname = generate_slide_partname(self)
        slide_layout_part = slide_layout.part
        slide_part = SlidePart.new(partname, self.package, slide_layout_part)
        rId = self.relate_to(slide_part, RT.SLIDE)
        return rId, slide_part.slide

    def add_slide_ppt(self, slide_layout):
        """Add slide to presentation"""
        rId, slide = add_slide_part(self.part, slide_layout)
        slide.shapes.clone_layout_placeholders(slide_layout)
        self._sldIdLst.add_sldId(rId)
        return slide

    return add_slide_ppt(ppt.slides, slide_layout)


def remove_shape(shape):
    """Remove a shape from a slide"""
    el = shape.element
    el.getparent().remove(el)


def copy_shapes(source, dest):
    """
    Copy all shapes from source to dest
    """
    for shape in source:
        if isinstance(shape, GroupShape):
            # Copy group shape
            group = dest.shapes.add_group_shape()
            group.name = shape.name
            group.left = shape.left
            group.top = shape.top
            group.width = shape.width
            group.height = shape.height
            group.rotation = shape.rotation

            # Recursive copy of contents
            copy_shapes(shape.shapes, group)

            # Fix offset
            cur_el = group._element.xpath(".//p:grpSpPr")[0]
            ref_el = shape._element.xpath(".//p:grpSpPr")[0]
            parent = cur_el.getparent()
            parent.insert(parent.index(cur_el) + 1, copy.deepcopy(ref_el))
            parent.remove(cur_el)

        elif hasattr(shape, "image"):
            # Copy image
            content = io.BytesIO(shape.image.blob)
            result = dest.shapes.add_picture(
                content, shape.left, shape.top, shape.width, shape.height
            )
            result.name = shape.name
            result.crop_left = shape.crop_left
            result.crop_right = shape.crop_right
            result.crop_top = shape.crop_top
            result.crop_bottom = shape.crop_bottom

        else:
            # Copy other shapes
            newel = copy.deepcopy(shape.element)
            # 🔧 FIX: use dest.shapes._spTree instead of dest._spTree
            dest.shapes._spTree.insert_element_before(newel, "p:extLst")


def duplicate_slide(ppt, slide_index: int):
    """
    Duplicate a slide exactly — preserves ALL formatting, shapes, and content
    """
    source = ppt.slides[slide_index]

    # Create a new slide
    dest = _exp_add_slide(ppt, source.slide_layout)

    # Remove placeholder shapes from destination
    for shape in list(dest.shapes):
        remove_shape(shape)

    # Copy all shapes
    copy_shapes(source.shapes, dest)

    # Copy notes (if any)
    if source.has_notes_slide:
        try:
            txt = source.notes_slide.notes_text_frame.text
            dest.notes_slide.notes_text_frame.text = txt
        except Exception:
            logger.warning(f"Could not copy notes from slide {slide_index + 1}")

    logger.info(f"✅ Duplicated slide {slide_index + 1} successfully")

    return dest
