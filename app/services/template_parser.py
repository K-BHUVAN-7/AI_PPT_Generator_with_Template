from pptx import Presentation
from typing import Dict, List, Tuple
import logging

logger = logging.getLogger(__name__)

class TemplateParser:
    """Parse and validate PowerPoint templates"""
    
    def __init__(self, template_path: str):
        self.template_path = template_path
        self.presentation = Presentation(template_path)
        self.layouts = {}
        self.placeholders_map = {}
        
    def parse_template(self) -> Dict:
        """Parse the template and extract layout information"""
        try:
            template_info = {
                "slide_layouts": [],
                "master_slides": [],
                "total_layouts": len(self.presentation.slide_layouts),
                "total_slides": len(self.presentation.slides)
            }
            
            # Parse each slide layout
            for idx, layout in enumerate(self.presentation.slide_layouts):
                layout_info = {
                    "index": idx,
                    "name": layout.name,
                    "placeholders": []
                }
                
                # Extract placeholder information
                for shape in layout.placeholders:
                    placeholder_info = {
                        "idx": shape.placeholder_format.idx,
                        "type": shape.placeholder_format.type,
                        "name": shape.name,
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height
                    }
                    layout_info["placeholders"].append(placeholder_info)
                
                template_info["slide_layouts"].append(layout_info)
                self.layouts[idx] = layout
            
            logger.info(f"Template parsed: {template_info['total_layouts']} layouts, {template_info['total_slides']} slides")
            return template_info
            
        except Exception as e:
            logger.error(f"Error parsing template: {str(e)}")
            raise
    
    def validate_template(self) -> Tuple[bool, List[str]]:
        """Validate that template meets minimum requirements"""
        errors = []
        
        # Check minimum slide count
        if len(self.presentation.slides) < 3:
            errors.append("Template must have at least 3 slides (cover, content, closing)")
        
        # Check for layouts
        if len(self.presentation.slide_layouts) < 1:
            errors.append("Template must have at least 1 slide layout")
        
        return len(errors) == 0, errors
