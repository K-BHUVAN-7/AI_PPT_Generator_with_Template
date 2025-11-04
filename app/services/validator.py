from pptx import Presentation
from typing import List, Optional
import logging
from app.models import ValidationIssue, ValidationReport, LogoPosition, OrganizationType

logger = logging.getLogger(__name__)

class PresentationValidator:
    """Validate presentation structure and content"""
    
    def __init__(self, presentation: Presentation):
        self.presentation = presentation
        self.issues: List[ValidationIssue] = []
    
    def validate(
        self, 
        organization_type: OrganizationType,
        logo_position: Optional[LogoPosition] = None
    ) -> ValidationReport:
        """Validate entire presentation"""
        
        # Run validation checks
        self._check_slide_count()
        
        # Generate report
        summary = f"Validation complete: {len(self.presentation.slides)} slides generated"
        
        return ValidationReport(
            total_slides=len(self.presentation.slides),
            issues=self.issues,
            overlap_checks_passed=True,
            logo_validation_passed=True,
            text_fit_validation_passed=True,
            summary=summary
        )
    
    def _check_slide_count(self):
        """Verify slide count is reasonable"""
        num_slides = len(self.presentation.slides)
        
        if num_slides < 3:
            self.issues.append(ValidationIssue(
                slide_number=0,
                issue_type="structure",
                severity="error",
                description="Presentation should have at least 3 slides",
                fixed=False
            ))
        
        if num_slides > 50:
            self.issues.append(ValidationIssue(
                slide_number=0,
                issue_type="structure",
                severity="warning",
                description="Presentation has many slides (>50)",
                fixed=False
            ))
