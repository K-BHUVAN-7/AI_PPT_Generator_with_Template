import aspose.slides as slides
from pptx import Presentation as PythonPptxPresentation
from typing import List, Optional, Dict
import logging
import os
import tempfile
from app.models import SlideContent

logger = logging.getLogger(__name__)

class SlideRenderer:
    """
    Smart template duplicator with configurable slide duplication
    """
    
    def __init__(self, template_path: str):
        self.template_path = os.path.abspath(template_path)
    
    async def render_presentation(
        self,
        slides_content: List[SlideContent],
        image_source: str = None,
        organization_type = None,
        logo_path: str = None,
        logo_position = None,
        template_config: Optional[Dict] = None  # NEW PARAMETER
    ) -> PythonPptxPresentation:
        """
        Duplicate template slides using custom configuration
        
        Args:
            slides_content: List of slide content (for count)
            template_config: Optional dict with unique_slides and duplicate_slides lists
        """
        
        template = slides.Presentation(self.template_path)
        num_template_slides = len(template.slides)
        num_required_slides = len(slides_content)
        
        logger.info(f"📊 Smart Template Duplication")
        logger.info(f"   Template slides: {num_template_slides}")
        logger.info(f"   Required slides: {num_required_slides}")
        
        if num_template_slides < 3:
            raise ValueError("Template must have at least 3 slides")
        
        # Use config or default
        if template_config:
            unique_slides = template_config.get('unique_slides', [])
            duplicate_slides = template_config.get('duplicate_slides', list(range(2, num_template_slides)))
            strategy = template_config.get('strategy', 'Custom')
        else:
            unique_slides = []
            duplicate_slides = list(range(2, num_template_slides))
            strategy = "Default: Duplicate all middle slides"
        
        logger.info(f"📋 Strategy: {strategy}")
        logger.info(f"   Unique slides (use once): {unique_slides}")
        logger.info(f"   Duplicate slides (repeat): {duplicate_slides}")
        
        # Start with cover slide
        output = slides.Presentation(self.template_path)
        
        while len(output.slides) > 1:
            output.slides.remove_at(1)
        
        logger.info(f"✓ Base: 1 slide (cover)")
        
        # Calculate content needed
        num_content_needed = num_required_slides - 2  # Exclude cover and closing
        
        # Add unique slides first (they appear only once)
        unique_added = 0
        for slide_num in unique_slides:
            if unique_added < num_content_needed:
                fresh_template = slides.Presentation(self.template_path)
                source_slide = fresh_template.slides[slide_num - 1]  # Convert to 0-based
                output.slides.add_clone(source_slide)
                unique_added += 1
                logger.info(f"   Added unique slide {slide_num}")
        
        # Add duplicate slides (cycle through them)
        remaining_needed = num_content_needed - unique_added
        
        if remaining_needed > 0 and duplicate_slides:
            logger.info(f"🔄 Duplicating {remaining_needed} content slides")
            
            for i in range(remaining_needed):
                slide_num = duplicate_slides[i % len(duplicate_slides)]
                
                fresh_template = slides.Presentation(self.template_path)
                source_slide = fresh_template.slides[slide_num - 1]  # Convert to 0-based
                output.slides.add_clone(source_slide)
                
                if (i + 1) % 10 == 0:
                    logger.info(f"   Duplicated {i + 1}/{remaining_needed} slides")
        
        logger.info(f"✓ Content: {num_content_needed} slides")
        
        # Add closing slide
        fresh_template = slides.Presentation(self.template_path)
        closing_slide = fresh_template.slides[num_template_slides - 1]
        output.slides.add_clone(closing_slide)
        
        logger.info(f"✓ Closing: 1 slide")
        logger.info(f"✅ Total: {len(output.slides)} slides")
        
        # Save and convert
        temp_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
        temp_path = temp_file.name
        temp_file.close()
        
        output.save(temp_path, slides.export.SaveFormat.PPTX)
        
        final_prs = PythonPptxPresentation(temp_path)
        
        # Remove watermarks
        logger.info(f"🧹 Removing watermarks...")
        final_prs = self._remove_watermarks(final_prs)
        
        try:
            os.unlink(temp_path)
        except:
            pass
        
        logger.info(f"🎉 Complete: {len(final_prs.slides)} slides (clean)")
        
        return final_prs
    
    def _remove_watermarks(self, presentation: PythonPptxPresentation) -> PythonPptxPresentation:
        """Remove Aspose watermarks"""
        
        watermark_keywords = [
            'evaluation only',
            'created with aspose',
            'aspose.slides',
            'aspose pty ltd',
            'copyright',
            'evaluation',
            'aspose',
            'trial version'
        ]
        
        total_removed = 0
        
        for slide in presentation.slides:
            shapes_to_remove = []
            
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    try:
                        text = shape.text_frame.text.lower()
                        
                        if any(kw in text for kw in watermark_keywords):
                            shapes_to_remove.append(shape)
                    except:
                        pass
            
            for shape in shapes_to_remove:
                try:
                    sp = shape.element
                    sp.getparent().remove(sp)
                    total_removed += 1
                except:
                    pass
        
        if total_removed > 0:
            logger.info(f"✓ Removed {total_removed} watermarks")
        
        return presentation
