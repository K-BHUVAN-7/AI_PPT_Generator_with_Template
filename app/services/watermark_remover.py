from pptx import Presentation
import logging

logger = logging.getLogger(__name__)

class WatermarkRemover:
    """
    Removes Aspose evaluation watermarks from presentations
    """
    
    @staticmethod
    def remove_watermarks(presentation: Presentation) -> Presentation:
        """
        Remove all watermark text from slides
        
        Args:
            presentation: Presentation object to clean
        
        Returns:
            Cleaned presentation
        """
        
        watermark_keywords = [
            'evaluation only',
            'created with aspose',
            'aspose.slides',
            'aspose pty ltd',
            'copyright 2004-2025aspose',
            'evaluation',
            'aspose'
        ]
        
        total_removed = 0
        
        for slide_idx, slide in enumerate(presentation.slides):
            shapes_to_remove = []
            
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    try:
                        text = shape.text_frame.text.lower()
                        
                        # Check if this shape contains watermark text
                        is_watermark = any(keyword in text for keyword in watermark_keywords)
                        
                        if is_watermark:
                            shapes_to_remove.append(shape)
                            logger.debug(f"   Found watermark in slide {slide_idx + 1}: '{text[:50]}'")
                    except:
                        pass
            
            # Remove watermark shapes
            for shape in shapes_to_remove:
                sp = shape.element
                sp.getparent().remove(sp)
                total_removed += 1
        
        if total_removed > 0:
            logger.info(f"✓ Removed {total_removed} watermark shapes")
        else:
            logger.info(f"✓ No watermarks found")
        
        return presentation
