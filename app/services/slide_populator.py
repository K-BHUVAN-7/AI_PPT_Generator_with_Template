from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from typing import List, Optional
import logging
import os
from datetime import datetime
from app.models import SlideContent, LogoPosition, OrganizationType
from app.services.image_service import ImageService

logger = logging.getLogger(__name__)

class SlidePopulator:
    """
    Populates slides - USES template placeholders when available
    """
    
    def __init__(self):
        self.image_service = ImageService()
    
    async def populate_presentation(
        self,
        presentation: Presentation,
        slides_content: List[SlideContent],
        image_source: str,
        organization_type: OrganizationType,
        logo_path: Optional[str] = None,
        logo_position: Optional[LogoPosition] = None
    ) -> Presentation:
        """Populate presentation"""
        
        logger.info(f"\n📝 Populating {len(presentation.slides)} slides...")
        
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height
        
        logger.info(f"Slide size: {slide_width/914400:.2f}\" x {slide_height/914400:.2f}\"")
        
        template_slide_count = len(presentation.slides)
        
        for idx, content in enumerate(slides_content):
            if idx >= len(presentation.slides):
                break
            
            slide = presentation.slides[idx]
            
            if idx == 0:
                logger.info(f"\n✏️  Slide {idx + 1} (COVER): '{content.title}'")
                self._add_title_slide_content_smart(presentation, slide, content)
            
            elif idx == template_slide_count - 1:
                logger.info(f"\n✏️  Slide {idx + 1} (CLOSING): Keeping as-is")
                continue
            
            else:
                logger.info(f"\n✏️  Slide {idx + 1} (CONTENT): '{content.title}'")
                await self._add_content_slide(presentation, slide, content, image_source)
        
        logger.info(f"\n✅ Population complete!")
        return presentation
    
    def _add_title_slide_content_smart(self, prs: Presentation, slide, content: SlideContent):
        """
        Add content to title slide - SMART detection of template placeholders
        """
        
        title = content.title
        subtitle = content.bullet_points[0] if content.bullet_points else ''
        
        # Find existing text shapes in template (sorted by size and position)
        text_shapes = []
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                try:
                    _ = shape.text_frame
                    # Get shape area (size)
                    area = shape.width * shape.height
                    text_shapes.append({
                        'shape': shape,
                        'area': area,
                        'top': shape.top,
                        'width': shape.width
                    })
                except:
                    pass
        
        # Sort by: 1) vertical position (top first), 2) size (larger first)
        text_shapes.sort(key=lambda s: (s['top'], -s['area']))
        
        logger.info(f"   Found {len(text_shapes)} text shapes in template")
        
        if len(text_shapes) >= 2:
            # USE TEMPLATE PLACEHOLDERS
            # First shape = title
            title_shape = text_shapes[0]['shape']
            title_frame = title_shape.text_frame
            title_frame.clear()
            
            p = title_frame.add_paragraph()
            p.text = title
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.name = 'Arial'
            p.alignment = PP_ALIGN.CENTER
            
            logger.info(f"   ✓ Used template title placeholder at ({title_shape.left/914400:.2f}\", {title_shape.top/914400:.2f}\")")
            
            # Second shape = subtitle
            if subtitle and len(text_shapes) >= 2:
                subtitle_shape = text_shapes[1]['shape']
                subtitle_frame = subtitle_shape.text_frame
                subtitle_frame.clear()
                
                p = subtitle_frame.add_paragraph()
                p.text = subtitle
                p.font.size = Pt(20)
                p.font.name = 'Arial'
                p.font.color.rgb = RGBColor(80, 80, 80)
                p.alignment = PP_ALIGN.CENTER
                
                logger.info(f"   ✓ Used template subtitle placeholder")
        
        else:
            # FALLBACK: Add new textboxes
            logger.warning(f"   ⚠ No template placeholders found, using fallback positioning")
            
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            margin = Inches(0.75)
            
            # Title
            title_left = margin
            title_width = slide_width - (2 * margin)
            title_top = int(slide_height * 0.35)
            title_height = Inches(1.2)
            
            title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.word_wrap = True
            title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(40)
            title_para.font.bold = True
            title_para.font.name = 'Arial'
            title_para.alignment = PP_ALIGN.CENTER
            
            # Subtitle
            if subtitle:
                subtitle_top = int(slide_height * 0.52)
                subtitle_height = Inches(0.8)
                
                subtitle_box = slide.shapes.add_textbox(title_left, subtitle_top, title_width, subtitle_height)
                subtitle_frame = subtitle_box.text_frame
                subtitle_frame.text = subtitle
                subtitle_frame.word_wrap = True
                
                subtitle_para = subtitle_frame.paragraphs[0]
                subtitle_para.font.size = Pt(20)
                subtitle_para.font.name = 'Arial'
                subtitle_para.font.color.rgb = RGBColor(80, 80, 80)
                subtitle_para.alignment = PP_ALIGN.CENTER
        
        logger.info(f"   ✓ Title slide complete")
    
    async def _add_content_slide(self, prs: Presentation, slide, content: SlideContent, image_source: str):
        """Add content to content slides - PREVENT OVERFLOW"""
    
        title = content.title
        bullets = content.bullet_points
    
        slide_width = prs.slide_width
        slide_height = prs.slide_height
    
        # INCREASED MARGINS - More conservative
        margin_h = Inches(0.7)  # INCREASED from 0.6 to 0.7
        margin_top = Inches(0.9)     # INCREASED from 0.8 to 0.9
        margin_bottom = Inches(0.5)  # INCREASED from 0.6 to 0.5
    
        usable_width = slide_width - (2 * margin_h)
        usable_height = slide_height - margin_top - margin_bottom
    
        # TITLE - Smaller height to give more space to content
        title_left = margin_h
        title_top = margin_top
        title_width = usable_width
        title_height = Inches(0.7)  # REDUCED from 0.8 to 0.7
    
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True
        title_frame.margin_bottom = 0
    
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)  # REDUCED from 28 to 26
        title_para.font.bold = True
        title_para.font.name = 'Arial'
        title_para.font.color.rgb = RGBColor(31, 73, 125)
        title_para.alignment = PP_ALIGN.LEFT
        title_para.space_after = Pt(0)
    
        logger.info(f"   ✓ Title: {title}")
    
        # Layout
        has_image = content.image_concept and image_source != "none"
    
        gap_after_title = Inches(0.2)  # REDUCED from 0.25 to 0.2
        content_top = title_top + title_height + gap_after_title
        available_height = slide_height - content_top - margin_bottom
    
        # Log available space
        logger.info(f"   Available text height: {available_height/914400:.2f}\"")
    
        if has_image:
            text_width = int(usable_width * 0.50)
            gap = int(usable_width * 0.05)
            image_width = int(usable_width * 0.45)
        
            text_left = margin_h
            image_left = text_left + text_width + gap
        
            if image_left + image_width > slide_width - margin_h:
                image_width = (slide_width - margin_h) - image_left
        else:
            text_left = margin_h
            text_width = usable_width
    
        # LIMIT BULLETS if too many
        max_bullets = 7
        if len(bullets) > max_bullets:
            logger.warning(f"   ⚠️  Too many bullets ({len(bullets)}), truncating to {max_bullets}")
            bullets = bullets[:max_bullets]
    
    # TEXT AREA
        text_box = slide.shapes.add_textbox(text_left, content_top, text_width, available_height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        text_frame.margin_left = Inches(0.15)
        text_frame.margin_right = Inches(0.15)
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)
    
        # Calculate font size based on number of bullets
        if len(bullets) <= 4:
            font_size = Pt(17)
            space_before = Pt(9)
            space_after = Pt(10)
        elif len(bullets) <= 6:
            font_size = Pt(15)  # Current
            space_before = Pt(7)
            space_after = Pt(8)
        else:  # 7+ bullets
            font_size = Pt(12)  # SMALLER
            space_before = Pt(4)
            space_after = Pt(6)
    
        logger.info(f"   Font size: {font_size.pt}pt for {len(bullets)} bullets")
    
        # Add bullets
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
        
            # TRUNCATE VERY LONG BULLETS
            max_bullet_length = 120
            if len(bullet) > max_bullet_length:
                bullet = bullet[:max_bullet_length-3] + "..."
                logger.warning(f"   ⚠️  Truncated bullet {i+1} to {max_bullet_length} chars")
        
            p.text = " " + bullet
            p.level = 0
        
            self._add_bullet_formatting(p)
        
            p.font.size = font_size
            p.font.name = 'Arial'
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.space_before = space_before
            p.space_after = space_after
            p.line_spacing = 1.1  # TIGHTER from 1.15
            p.alignment = PP_ALIGN.LEFT

        logger.info(f"   ✓ Bullets: {len(bullets)}")
    
        # Add image
        if has_image:
            await self._add_image_strict(
                prs, slide, content.image_concept, image_source,
                image_left, content_top, image_width, available_height
            )
    
    # Notes
        if content.speaker_notes:
            try:
                slide.notes_slide.notes_text_frame.text = content.speaker_notes
                logger.info(f"   ✓ Notes")
            except:
                pass

    
    def _add_bullet_formatting(self, paragraph):
        """Add bullet formatting"""
        pPr = paragraph._element.get_or_add_pPr()
        
        buFont = OxmlElement('a:buFont')
        buFont.set('typeface', 'Arial')
        
        buChar = OxmlElement('a:buChar')
        buChar.set('char', '•')
        
        pPr.insert(0, buFont)
        pPr.insert(1, buChar)
    
    async def _add_image_strict(
        self, prs: Presentation, slide, image_prompt: str, image_source: str,
        left: int, top: int, max_width: int, max_height: int
    ):
        """Add image with strict bounds"""
        
        try:
            logger.info(f"   📷 Searching: '{image_prompt}'")
            
            image_url = self.image_service.get_stock_image(image_prompt)
            
            if not image_url:
                logger.warning(f"   ⚠ No image")
                return
            
            temp_path = f"temp_image_{datetime.now().timestamp()}.jpg"
            
            if not self.image_service.download_image(image_url, temp_path):
                logger.warning(f"   ⚠ Download failed")
                return
            
            try:
                from PIL import Image
                
                img = Image.open(temp_path)
                img_width, img_height = img.size
                aspect = img_width / img_height
                
                if max_width / aspect <= max_height:
                    width = max_width
                    height = int(max_width / aspect)
                else:
                    height = max_height
                    width = int(max_height * aspect)
                
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                
                if left + width > slide_width:
                    width = slide_width - left
                    height = int(width / aspect)
                
                if top + height > slide_height:
                    height = slide_height - top
                    width = int(height * aspect)
                
                if left < 0 or top < 0 or left + width > slide_width or top + height > slide_height:
                    logger.error(f"   ✗ Invalid bounds")
                    return
                
                slide.shapes.add_picture(temp_path, left, top, width=width, height=height)
                
                logger.info(f"   ✅ Image ({width/914400:.2f}\" x {height/914400:.2f}\")")
            
            finally:
                if os.path.exists(temp_path):
                    os.remove(temp_path)
        
        except Exception as e:
            logger.error(f"   ✗ Image error: {str(e)}")
